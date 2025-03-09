"""
Slide class for the pptx_charts_tables package.
This module provides the PPTXSlide class for managing slide content.
"""

from pptx.util import Inches
from .tables.styled_table import StyledTable
from .charts.bar_chart import BarChart
from .charts.clustered_bar import ClusteredBarChart
from .charts.stacked_bar import StackedBarChart
from .charts.donut_chart import DonutChart
from .charts.line_chart import LineChart
from .utils.shapes import add_text_box
from .utils.shapes import add_shape
from .utils.shapes import add_arrow
from .utils.shapes import add_image
from .utils.conditional_formatting import apply_conditional_formatting
from .utils.colors import create_palette
from .utils.colors import get_color_scheme


class PPTXSlide:
    """
    Class for managing slide content including charts and tables.
    """

    def __init__(self, presentation, slide, config):
        """
        Initialize a slide.

        Args:
            presentation (PPTXPresentation): Parent presentation.
            slide (pptx.slide.Slide): The underlying slide object.
            config (dict): Configuration for this slide.
        """
        self.presentation = presentation
        self.slide = slide
        self.config = config
        self.charts = []
        self.tables = []
        self.shapes = []

    def add_chart(self, chart_type, data, position, size, **kwargs):
        """
        Add a chart to the slide.

        Args:
            chart_type (str): Type of chart ('bar', 'clustered_bar', 'stacked_bar', 'donut', 'line').
            data (pd.DataFrame): Data for the chart.
            position (tuple): (x, y) position in inches.
            size (tuple): (width, height) size in inches.
            **kwargs: Additional chart-specific options.

        Returns:
            Chart: The created chart object.
        """

        x, y = position
        width, height = size

        chart_classes = {
            "bar": BarChart,
            "clustered_bar": ClusteredBarChart,
            "stacked_bar": StackedBarChart,
            "donut": DonutChart,
            "line": LineChart,
        }

        if chart_type not in chart_classes:
            raise ValueError(f"Unsupported chart type: {chart_type}")

        chart_class = chart_classes[chart_type]
        chart = chart_class(
            self,
            data,
            (Inches(x), Inches(y)),
            (Inches(width), Inches(height)),
            self.config,
            **kwargs,
        )

        self.charts.append(chart)
        return chart

    def add_table(self, data, position, **kwargs):
        """
        Add a table to the slide.

        Args:
            data (pd.DataFrame): Data for the table.
            position (tuple): (x, y) position in inches.
            **kwargs: Additional table-specific options.

        Returns:
            StyledTable: The created table object.
        """

        x, y = position

        table = StyledTable(self, data, (Inches(x), Inches(y)), self.config, **kwargs)

        self.tables.append(table)
        return table

    def add_text_box(self, text, position, size=None, **kwargs):
        """
        Add a text box to the slide.

        Args:
            text (str): Text content
            position (tuple): (x, y) position in inches
            size (tuple, optional): (width, height) size in inches
            **kwargs: Additional style options

        Returns:
            pptx.shapes.Shape: The created text box shape
        """

        text_box = add_text_box(self, text, position, size, **kwargs)
        self.shapes.append(text_box)
        return text_box

    def add_shape(self, shape_type, position, size, **kwargs):
        """
        Add a shape to the slide.

        Args:
            shape_type (str): Shape type ('rectangle', 'oval', etc.)
            position (tuple): (x, y) position in inches
            size (tuple): (width, height) size in inches
            **kwargs: Additional style options

        Returns:
            pptx.shapes.Shape: The created shape
        """

        shape = add_shape(self, shape_type, position, size, **kwargs)
        self.shapes.append(shape)
        return shape

    def add_arrow(self, start_pos, end_pos, **kwargs):
        """
        Add an arrow connecting two points.

        Args:
            start_pos (tuple): (x, y) starting position in inches
            end_pos (tuple): (x, y) ending position in inches
            **kwargs: Additional style options

        Returns:
            pptx.shapes.Shape: The created connector shape
        """

        arrow = add_arrow(self, start_pos, end_pos, **kwargs)
        self.shapes.append(arrow)
        return arrow

    def add_image(self, image_path, position, size=None, **kwargs):
        """
        Add an image to the slide.

        Args:
            image_path (str): Path to the image file
            position (tuple): (x, y) position in inches
            size (tuple, optional): (width, height) size in inches
            **kwargs: Additional options

        Returns:
            pptx.shapes.Picture: The created picture shape
        """

        image = add_image(self, image_path, position, size, **kwargs)
        self.shapes.append(image)
        return image

    def apply_conditional_formatting(self, table, rules, start_row=0, **kwargs):
        """
        Apply conditional formatting to a table.

        Args:
            table: Table object (from add_table)
            rules (list): List of formatting rules
            start_row (int): Row to start applying formatting
            **kwargs: Additional options
        """

        # If table is an index or a StyledTable object, get the actual table
        if isinstance(table, int) and table < len(self.tables):
            table = self.tables[table].table
        elif hasattr(table, "table"):
            table = table.table

        apply_conditional_formatting(table, rules, start_row, **kwargs)

    def create_color_palette(self, base_color, variations=5, mode="monochromatic"):
        """
        Create a color palette based on a base color.

        Args:
            base_color (str): Base hex color
            variations (int): Number of variations
            mode (str): Type of palette - 'monochromatic', 'complementary', 'analogous'

        Returns:
            list: List of hex color strings
        """

        return create_palette(base_color, variations, mode)

    def get_color_scheme(self, scheme_name):
        """
        Get a predefined color scheme.

        Args:
            scheme_name (str): Name of the color scheme

        Returns:
            list: List of hex color strings
        """

        return get_color_scheme(scheme_name)
