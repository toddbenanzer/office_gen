"""
Base chart class for the pptx_charts_tables package.
This module provides the Chart base class from which all chart types inherit.
"""

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK


class Chart:
    """
    Base class for all chart types.
    """

    def __init__(self, slide, data, position, size, config, **kwargs):
        """
        Initialize a chart.

        Args:
            slide (PPTXSlide): Parent slide.
            data (pd.DataFrame): Data for the chart.
            position (tuple): (x, y) position in inches.
            size (tuple): (width, height) size in inches.
            config (dict): Configuration.
            **kwargs: Additional chart-specific options.
        """
        self.slide = slide
        self.data = data
        self.position = position
        self.size = size
        self.config = config
        self.kwargs = kwargs
        self.chart = None

    def _create_chart(self):
        """Create the chart. To be implemented by subclasses."""
        raise NotImplementedError

    def _apply_styles(self):
        """Apply styles to the chart. To be implemented by subclasses."""
        raise NotImplementedError

    def _set_chart_title(self, title=None):
        """
        Set the chart title.

        Args:
            title (str, optional): The title text. If None, use from kwargs.
        """
        title_text = title or self.kwargs.get("title")
        if not title_text:
            self.chart.has_title = False
            return

        self.chart.has_title = True
        self.chart.chart_title.has_text_frame = True
        self.chart.chart_title.text_frame.text = title_text

        font_name = self.kwargs.get(
            "title_font_name", self.config["general"]["font_name"]
        )
        font_size = self.kwargs.get(
            "title_font_size", self.config["general"]["title_font_size"]
        )

        self.chart.chart_title.text_frame.paragraphs[0].font.name = font_name
        self.chart.chart_title.text_frame.paragraphs[0].font.size = Pt(font_size)

    def _setup_category_axis(self):
        """Set up the category axis."""
        if not hasattr(self.chart, "category_axis"):
            return

        category_axis = self.chart.category_axis

        # Gridlines
        category_axis.has_major_gridlines = self.kwargs.get(
            "category_axis_has_gridlines",
            self.config["charts"].get("category_axis_has_gridlines", False),
        )

        # Tick marks
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.minor_tick_mark = XL_TICK_MARK.NONE

        # Font
        font_name = self.kwargs.get(
            "category_axis_font_name",
            self.config["charts"].get(
                "axis_font_name", self.config["general"]["font_name"]
            ),
        )
        font_size = self.kwargs.get(
            "category_axis_font_size",
            self.config["charts"].get(
                "axis_font_size", self.config["general"]["font_size"]
            ),
        )

        category_axis.tick_labels.font.name = font_name
        category_axis.tick_labels.font.size = Pt(font_size)

        # Remove axis line
        if not self.kwargs.get(
            "show_axis_lines", self.config["charts"].get("show_axis_lines", False)
        ):
            category_axis.format.line.fill.background()

    def _setup_value_axis(self):
        """Set up the value axis."""
        if not hasattr(self.chart, "value_axis"):
            return

        value_axis = self.chart.value_axis

        # Visibility
        value_axis.visible = self.kwargs.get(
            "value_axis_visible", self.config["charts"].get("value_axis_visible", False)
        )

        # Gridlines
        value_axis.has_major_gridlines = self.kwargs.get(
            "value_axis_has_gridlines",
            self.config["charts"].get("value_axis_has_gridlines", False),
        )
        value_axis.has_minor_gridlines = self.kwargs.get(
            "value_axis_has_minor_gridlines",
            self.config["charts"].get("value_axis_has_minor_gridlines", False),
        )

        # Tick marks
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.minor_tick_mark = XL_TICK_MARK.NONE

        # Number format
        if "value_axis_number_format" in self.kwargs:
            value_axis.tick_labels.number_format = self.kwargs[
                "value_axis_number_format"
            ]

        # Font
        if value_axis.visible:
            font_name = self.kwargs.get(
                "value_axis_font_name",
                self.config["charts"].get(
                    "axis_font_name", self.config["general"]["font_name"]
                ),
            )
            font_size = self.kwargs.get(
                "value_axis_font_size",
                self.config["charts"].get(
                    "axis_font_size", self.config["general"]["font_size"]
                ),
            )

            value_axis.tick_labels.font.name = font_name
            value_axis.tick_labels.font.size = Pt(font_size)

    def _setup_legend(self):
        """Set up the chart legend."""
        self.chart.has_legend = self.kwargs.get(
            "has_legend", self.config["charts"].get("has_legend", True)
        )

        if not self.chart.has_legend:
            return

        legend = self.chart.legend

        # Font
        font_name = self.kwargs.get(
            "legend_font_name",
            self.config["charts"].get(
                "legend_font_name", self.config["general"]["font_name"]
            ),
        )
        font_size = self.kwargs.get(
            "legend_font_size",
            self.config["charts"].get(
                "legend_font_size", self.config["general"]["font_size"]
            ),
        )

        legend.font.name = font_name
        legend.font.size = Pt(font_size)

        # Position (if supported by python-pptx version)
        position = self.kwargs.get(
            "legend_position", self.config["charts"].get("legend_position", "bottom")
        )

        if hasattr(legend, "position"):
            from pptx.enum.chart import XL_LEGEND_POSITION

            positions = {
                "top": XL_LEGEND_POSITION.TOP,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
                "corner": XL_LEGEND_POSITION.CORNER,
            }
            if position in positions:
                legend.position = positions[position]

    def _setup_data_labels(self):
        """Set up data labels for the chart."""
        if not hasattr(self.chart, "plots") or not self.chart.plots:
            return

        plot = self.chart.plots[0]

        has_data_labels = self.kwargs.get(
            "has_data_labels", self.config["charts"].get("has_data_labels", True)
        )

        plot.has_data_labels = has_data_labels

        if not has_data_labels:
            return

        data_labels = plot.data_labels

        # Font
        font_name = self.kwargs.get(
            "data_label_font_name",
            self.config["charts"].get(
                "data_label_font_name", self.config["general"]["font_name"]
            ),
        )
        font_size = self.kwargs.get(
            "data_label_font_size",
            self.config["charts"].get(
                "data_label_font_size", self.config["general"]["font_size"]
            ),
        )

        data_labels.font.name = font_name
        data_labels.font.size = Pt(font_size)

        # Number format
        if "data_label_number_format" in self.kwargs:
            data_labels.number_format = self.kwargs["data_label_number_format"]

        # Position (if supported)
        if "data_label_position" in self.kwargs:
            from pptx.enum.chart import XL_LABEL_POSITION

            positions = {
                "inside_end": XL_LABEL_POSITION.INSIDE_END,
                "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
                "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
                "center": XL_LABEL_POSITION.CENTER,
            }
            position = self.kwargs["data_label_position"]
            if position in positions:
                data_labels.position = positions[position]
