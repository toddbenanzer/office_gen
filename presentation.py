"""
Main presentation class for the pptx_charts_tables package.
This module provides the PPTXPresentation class, which is the main entry point
for creating and managing PowerPoint presentations.
"""

from pptx import Presentation
import copy
from .config import DEFAULT_CONFIG
from .slide import PPTXSlide


class PPTXPresentation:
    """
    Main class for creating and managing PowerPoint presentations.
    """

    def __init__(self, template=None, config=None):
        """
        Initialize a new presentation.

        Args:
            template (str, optional): Path to a template file.
            config (dict, optional): Custom configuration to override defaults.
        """
        self.config = copy.deepcopy(DEFAULT_CONFIG)
        if config:
            self._update_config(config)

        if template:
            self.presentation = Presentation(template)
        else:
            self.presentation = Presentation()

        self.slides = []

    def _update_config(self, config):
        """
        Update the default configuration with custom settings.

        Args:
            config (dict): Custom configuration settings.
        """
        for section, values in config.items():
            if section in self.config:
                if isinstance(values, dict) and isinstance(self.config[section], dict):
                    self.config[section].update(values)
                else:
                    self.config[section] = values
            else:
                self.config[section] = values

    def add_slide(self, layout_type=None, title=None, config=None):
        """
        Add a new slide to the presentation.

        Args:
            layout_type (int, optional): Layout type index. Default is 5 (Title Only).
            title (str, optional): Slide title.
            config (dict, optional): Custom configuration for this slide.

        Returns:
            PPTXSlide: The newly created slide.
        """
        if layout_type is None:
            layout_type = 5  # Title Only layout by default

        slide_layout = self.presentation.slide_layouts[layout_type]
        slide = self.presentation.slides.add_slide(slide_layout)

        if title and hasattr(slide.shapes, "title") and slide.shapes.title:
            slide.shapes.title.text = title

        slide_config = copy.deepcopy(self.config)
        if config:
            self._update_config_for_slide(slide_config, config)

        pptx_slide = PPTXSlide(self, slide, slide_config)
        self.slides.append(pptx_slide)
        return pptx_slide

    def _update_config_for_slide(self, slide_config, config):
        """
        Update the configuration for a specific slide.

        Args:
            slide_config (dict): Current slide configuration.
            config (dict): New configuration values.
        """
        for section, values in config.items():
            if section in slide_config:
                if isinstance(values, dict) and isinstance(slide_config[section], dict):
                    slide_config[section].update(values)
                else:
                    slide_config[section] = values
            else:
                slide_config[section] = values

    def save(self, path):
        """
        Save the presentation to a file.

        Args:
            path (str): Path to save the presentation.
        """
        self.presentation.save(path)
