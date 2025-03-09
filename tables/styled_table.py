"""
Styled table implementation for the pptx_charts_tables package.
"""

from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import pandas as pd
from ..utils.formatting import format_value


class StyledTable:
    """
    Styled table implementation.
    """

    def __init__(self, slide, data, position, config, **kwargs):
        """
        Initialize a styled table.

        Args:
            slide (PPTXSlide): Parent slide.
            data (pd.DataFrame): Data for the table.
            position (tuple): (x, y) position in inches.
            config (dict): Configuration.
            **kwargs: Additional table-specific options such as:
                - col_widths (list): Column widths in inches.
                - has_header (bool): Whether the first row is a header.
                - header_style (dict): Style for header row.
                - column_formats (dict): Formatting for specific columns.
                - row_styles (list): Styles for specific rows.
                - merged_cells (list): List of cell ranges to merge.
                - total_rows (list): Indices of total rows for special formatting.
                - subtotal_rows (list): Indices of subtotal rows for special formatting.
                - alternating_row_fill (bool): Whether to use alternating row fill.
        """
        self.slide = slide
        self.data = data
        self.position = position
        self.config = config
        self.kwargs = kwargs
        self.table = None

        self._create_table()
        self._apply_styles()

    def _create_table(self):
        """Create the table in the slide."""
        # Calculate table dimensions
        rows, cols = self.data.shape

        # Add header row if needed
        has_header = self.kwargs.get(
            "has_header", self.config["tables"].get("has_header", True)
        )

        if has_header:
            table_rows = rows + 1
        else:
            table_rows = rows

        # Calculate column widths
        col_widths = self.kwargs.get("col_widths")
        if not col_widths:
            # Default to equal widths
            total_width = self.kwargs.get("table_width", 8)  # default 8 inches
            col_widths = [total_width / cols] * cols

        # Get row height
        row_height = self.kwargs.get(
            "row_height",
            self.config["tables"].get("row_height", 0.3),  # default 0.3 inches
        )

        # Create table
        table = self.slide.slide.shapes.add_table(
            table_rows,
            cols,
            self.position[0],
            self.position[1],
            Inches(sum(col_widths)),
            Inches(table_rows * row_height),
        ).table

        # Set column widths
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

        # Populate table
        if has_header:
            # Add header row
            for i, col in enumerate(self.data.columns):
                table.cell(0, i).text = str(col)

            # Add data rows
            for row_idx, (_, row) in enumerate(self.data.iterrows(), start=1):
                for col_idx, value in enumerate(row):
                    table.cell(row_idx, col_idx).text = self._format_cell_value(
                        value, col_idx
                    )
        else:
            # Add data rows without header
            for row_idx, (_, row) in enumerate(self.data.iterrows()):
                for col_idx, value in enumerate(row):
                    table.cell(row_idx, col_idx).text = self._format_cell_value(
                        value, col_idx
                    )

        self.table = table

    def _format_cell_value(self, value, col_idx):
        """
        Format cell value based on column format.

        Args:
            value: The cell value.
            col_idx (int): Column index.

        Returns:
            str: Formatted value as string.
        """
        # Get column name
        col_name = self.data.columns[col_idx]

        # Get format type for this column
        column_formats = self.kwargs.get("column_formats", {})
        format_type = column_formats.get(col_name, None)

        if format_type is None:
            # Try to infer format type from column name
            col_lower = col_name.lower()
            if any(
                x in col_lower
                for x in ["price", "cost", "revenue", "dollar", "$", "sales"]
            ):
                format_type = "dollars"
            elif any(x in col_lower for x in ["percent", "%", "rate", "ratio"]):
                format_type = "percentage"
            elif any(
                x in col_lower for x in ["count", "number", "total", "qty", "quantity"]
            ):
                format_type = "counts"
            else:
                format_type = "text"

        # Format based on type
        return format_value(value, format_type, self.config)

    def _apply_styles(self):
        """Apply styles to the table."""
        # Apply header styles
        self._apply_header_styles()

        # Apply alternating row fill
        self._apply_alternating_row_fill()

        # Apply specific row styles
        self._apply_row_styles()

        # Apply specific column styles
        self._apply_column_styles()

        # Apply merged cells
        self._apply_merged_cells()

        # Apply special formatting for total and subtotal rows
        self._apply_total_row_styles()

    def _apply_header_styles(self):
        """Apply styles to the header row."""
        has_header = self.kwargs.get(
            "has_header", self.config["tables"].get("has_header", True)
        )

        if not has_header:
            return

        header_style = self.kwargs.get("header_style", {})

        # Get header styles from config
        font_name = header_style.get(
            "font_name",
            self.config["tables"].get(
                "header_font_name", self.config["general"]["font_name"]
            ),
        )
        font_size = header_style.get(
            "font_size",
            self.config["tables"].get(
                "header_font_size", self.config["general"]["font_size"]
            ),
        )
        font_bold = header_style.get(
            "font_bold", self.config["tables"].get("header_font_bold", True)
        )
        fill_color = header_style.get(
            "fill_color", self.config["tables"].get("header_fill_color", "3C2F80")
        )
        font_color = header_style.get(
            "font_color", self.config["tables"].get("header_font_color", "FFFFFF")
        )

        # Apply styles to header row
        for i in range(len(self.data.columns)):
            cell = self.table.cell(0, i)

            # Apply text styles
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = Pt(font_size)
                paragraph.font.bold = font_bold
                paragraph.font.color.rgb = RGBColor.from_string(font_color)
                paragraph.alignment = PP_ALIGN.CENTER

            # Apply cell styles
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)

            # Apply alignment
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _apply_alternating_row_fill(self):
        """Apply alternating row fill colors."""
        alternating_row_fill = self.kwargs.get(
            "alternating_row_fill",
            self.config["tables"].get("alternating_row_fill", True),
        )

        if not alternating_row_fill:
            return

        alternating_row_fill_color = self.kwargs.get(
            "alternating_row_fill_color",
            self.config["tables"].get("alternating_row_fill_color", "F2F2F2"),
        )

        has_header = self.kwargs.get(
            "has_header", self.config["tables"].get("has_header", True)
        )

        start_row = 1 if has_header else 0

        for row_idx in range(start_row, len(self.table.rows)):
            # Apply alternating fill to even rows (0-indexed, but if header exists, start at 1)
            if (row_idx - start_row) % 2 == 1:
                for col_idx in range(len(self.table.columns)):
                    cell = self.table.cell(row_idx, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        alternating_row_fill_color
                    )

    def _apply_row_styles(self):
        """Apply styles to specific rows."""
        row_styles = self.kwargs.get("row_styles", [])

        for row_style in row_styles:
            row_idx = row_style.get("row_idx")
            if row_idx is None:
                continue

            fill_color = row_style.get("fill_color")
            font_name = row_style.get("font_name")
            font_size = row_style.get("font_size")
            font_bold = row_style.get("font_bold")
            font_color = row_style.get("font_color")

            for col_idx in range(len(self.table.columns)):
                cell = self.table.cell(row_idx, col_idx)

                if fill_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)

                if font_name or font_size or font_bold is not None or font_color:
                    for paragraph in cell.text_frame.paragraphs:
                        if font_name:
                            paragraph.font.name = font_name
                        if font_size:
                            paragraph.font.size = Pt(font_size)
                        if font_bold is not None:
                            paragraph.font.bold = font_bold
                        if font_color:
                            paragraph.font.color.rgb = RGBColor.from_string(font_color)

    def _apply_column_styles(self):
        """Apply styles to specific columns."""
        column_styles = self.kwargs.get("column_styles", [])

        for col_style in column_styles:
            col_idx = col_style.get("col_idx")
            if col_idx is None:
                continue

            alignment = col_style.get("alignment")

            has_header = self.kwargs.get(
                "has_header", self.config["tables"].get("has_header", True)
            )

            start_row = 1 if has_header else 0

            for row_idx in range(start_row, len(self.table.rows)):
                cell = self.table.cell(row_idx, col_idx)

                if alignment:
                    for paragraph in cell.text_frame.paragraphs:
                        if alignment == "left":
                            paragraph.alignment = PP_ALIGN.LEFT
                        elif alignment == "center":
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif alignment == "right":
                            paragraph.alignment = PP_ALIGN.RIGHT

    def _apply_merged_cells(self):
        """Apply merged cells."""
        merged_cells = self.kwargs.get("merged_cells", [])

        for merge in merged_cells:
            start_row = merge.get("start_row", 0)
            end_row = merge.get("end_row", start_row)
            start_col = merge.get("start_col", 0)
            end_col = merge.get("end_col", start_col)

            # Merge cells - python-pptx requires merging adjacent cells
            # First merge horizontally in each row
            for row in range(start_row, end_row + 1):
                for col in range(start_col, end_col):
                    try:
                        self.table.cell(row, col).merge(self.table.cell(row, col + 1))
                    except:
                        # Handle case where cells are already merged
                        pass

            # Then merge vertically for each column
            for col in range(start_col, end_col + 1):
                for row in range(start_row, end_row):
                    try:
                        self.table.cell(row, col).merge(self.table.cell(row + 1, col))
                    except:
                        # Handle case where cells are already merged
                        pass

    def _apply_total_row_styles(self):
        """Apply special formatting for total and subtotal rows."""
        # Total rows
        total_rows = self.kwargs.get("total_rows", [])

        for row_idx in total_rows:
            for col_idx in range(len(self.table.columns)):
                cell = self.table.cell(row_idx, col_idx)

                # Apply total row styles - bold font
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = self.config["tables"].get(
                        "totals_font_bold", True
                    )

        # Subtotal rows
        subtotal_rows = self.kwargs.get("subtotal_rows", [])

        for row_idx in subtotal_rows:
            for col_idx in range(len(self.table.columns)):
                cell = self.table.cell(row_idx, col_idx)

                # Apply subtotal row styles - bold font
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
