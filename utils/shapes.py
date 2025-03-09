"""
Shape utilities for the pptx_charts_tables package.
Provides functions for working with shapes, text boxes, arrows, and other elements.
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from ..config import DEFAULT_CONFIG


def add_text_box(slide, text, position, size=None, **kwargs):
    """
    Add a text box to a slide.

    Args:
        slide (PPTXSlide): The slide object
        text (str): Text content
        position (tuple): (x, y) position in inches
        size (tuple, optional): (width, height) size in inches
        **kwargs: Additional style options (font_name, font_size, bold, italic,
                 color, fill_color, align, v_align, etc.)

    Returns:
        pptx.shapes.Shape: The created text box shape
    """
    x, y = position

    # Default size if not specified
    if size is None:
        width = Inches(DEFAULT_CONFIG['text_box']['size']['width'])
        height = Inches(DEFAULT_CONFIG['text_box']['size']['height'])
    else:
        width, height = Inches(size[0]), Inches(size[1])
    # Default options if not specified
    if "align" not in kwargs:
        kwargs["align"] = DEFAULT_CONFIG['text_box']['alignment']['horizontal']
    if "v_align" not in kwargs:
        kwargs["v_align"] = DEFAULT_CONFIG['text_box']['alignment']['vertical']
    if "font_name" not in kwargs:
        kwargs["font_name"] = DEFAULT_CONFIG['text_box']['font']['name']
    if "font_size" not in kwargs:
        kwargs["font_size"] = DEFAULT_CONFIG['text_box']['font']['size']
    if "bold" not in kwargs:
        kwargs["bold"] = DEFAULT_CONFIG['text_box']['font']['bold']
    if "italic" not in kwargs:
        kwargs["italic"] = DEFAULT_CONFIG['text_box']['font']['italic']
    if "color" not in kwargs:
        kwargs["color"] = DEFAULT_CONFIG['text_box']['font']['color']
    if "fill_color" not in kwargs:
        kwargs["fill_color"] = DEFAULT_CONFIG['text_box']['fill']['color']
    if "no_fill" not in kwargs:
        kwargs["no_fill"] = DEFAULT_CONFIG['text_box']['fill']['no_fill']
    if "border_color" not in kwargs:
        kwargs["border_color"] = DEFAULT_CONFIG['text_box']['border']['color']
    if "border_width" not in kwargs:
        kwargs["border_width"] = DEFAULT_CONFIG['text_box']['border']['weight']
    if "no_border" not in kwargs:
        kwargs["no_border"] = DEFAULT_CONFIG['text_box']['border']['no_border']
    


    # Create shape
    text_box = slide.slide.shapes.add_textbox(Inches(x), Inches(y), width, height)

    # Add text
    text_frame = text_box.text_frame
    text_frame.text = text

    # Apply formatting
    paragraph = text_frame.paragraphs[0]

    # Text alignment
    alignment = kwargs["align"].lower()
    if alignment == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    elif alignment == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif alignment == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    elif alignment == "justify":
        paragraph.alignment = PP_ALIGN.JUSTIFY

    # Vertical alignment
    v_alignment = kwargs["v_align"].lower()
    if v_alignment == "top":
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
    elif v_alignment == "middle":
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_alignment == "bottom":
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    # Font formatting
    font = paragraph.font

    if "font_name" in kwargs:
        font.name = kwargs["font_name"]

    if "font_size" in kwargs:
        font.size = Pt(kwargs["font_size"])
    else:
        font.size = Pt(DEFAULT_CONFIG['text_box']['font_size'])

    if "bold" in kwargs:
        font.bold = kwargs["bold"]

    if "italic" in kwargs:
        font.italic = kwargs["italic"]

    if "color" in kwargs:
        font.color.rgb = RGBColor.from_string(kwargs["color"])

    # Background color
    if "fill_color" in kwargs:
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(kwargs["fill_color"])

    # No fill
    if kwargs.get("no_fill", False):
        text_box.fill.background()

    # Border
    if "border_color" in kwargs:
        line = text_box.line
        line.color.rgb = RGBColor.from_string(kwargs["border_color"])

        # Line width (in points)
        if "border_width" in kwargs:
            line.width = Pt(kwargs["border_width"])

    # No border
    if kwargs.get("no_border", False):
        text_box.line.fill.background()

    return text_box


def add_shape(slide, shape_type, position, size, **kwargs):
    """
    Add a shape to a slide.

    Args:
        slide (PPTXSlide): The slide object
        shape_type (str): Shape type ('rectangle', 'oval', 'rounded_rectangle', etc.)
        position (tuple): (x, y) position in inches
        size (tuple): (width, height) size in inches
        **kwargs: Additional style options (fill_color, line_color, line_width, etc.)

    Returns:
        pptx.shapes.Shape: The created shape
    """
    x, y = position
    width, height = size

    # Map string shape types to pptx constants
    shape_types = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "chevron": MSO_SHAPE.CHEVRON,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "heptagon": MSO_SHAPE.HEPTAGON,
        "octagon": MSO_SHAPE.OCTAGON,
        "star": MSO_SHAPE.STAR_5_POINT,
        "cube": MSO_SHAPE.CUBE,
        "arc": MSO_SHAPE.ARC,
        "heart": MSO_SHAPE.HEART,
        "lightning": MSO_SHAPE.LIGHTNING_BOLT,
        "sun": MSO_SHAPE.SUN,
        "moon": MSO_SHAPE.MOON,
        "cloud": MSO_SHAPE.CLOUD,
        "smiley": MSO_SHAPE.SMILEY_FACE,
    }

    # Get the shape type constant
    mso_shape_type = shape_types.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)

    # Create shape
    shape = slide.slide.shapes.add_shape(
        mso_shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
    )

    # Fill color
    if "fill_color" in kwargs:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(kwargs["fill_color"])

    # No fill
    if kwargs.get("no_fill", False):
        shape.fill.background()

    # Line color and width
    if "line_color" in kwargs:
        shape.line.color.rgb = RGBColor.from_string(kwargs["line_color"])

        # Line width (in points)
        if "line_width" in kwargs:
            shape.line.width = Pt(kwargs["line_width"])

    # No line
    if kwargs.get("no_line", False):
        shape.line.fill.background()

    # Add text if provided
    if "text" in kwargs:
        text_frame = shape.text_frame
        text_frame.text = kwargs["text"]

        # Text alignment
        if "align" in kwargs:
            paragraph = text_frame.paragraphs[0]
            alignment = kwargs["align"].lower()
            if alignment == "left":
                paragraph.alignment = PP_ALIGN.LEFT
            elif alignment == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            elif alignment == "justify":
                paragraph.alignment = PP_ALIGN.JUSTIFY

        # Vertical alignment
        if "v_align" in kwargs:
            v_alignment = kwargs["v_align"].lower()
            if v_alignment == "top":
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
            elif v_alignment == "middle":
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif v_alignment == "bottom":
                text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        # Font formatting
        font = text_frame.paragraphs[0].font

        if "font_name" in kwargs:
            font.name = kwargs["font_name"]

        if "font_size" in kwargs:
            font.size = Pt(kwargs["font_size"])

        if "font_color" in kwargs:
            font.color.rgb = RGBColor.from_string(kwargs["font_color"])

        if "bold" in kwargs:
            font.bold = kwargs["bold"]

        if "italic" in kwargs:
            font.italic = kwargs["italic"]

    return shape


def add_arrow(slide, start_pos, end_pos, **kwargs):
    """
    Add an arrow connecting two points.

    Args:
        slide (PPTXSlide): The slide object
        start_pos (tuple): (x, y) starting position in inches
        end_pos (tuple): (x, y) ending position in inches
        **kwargs: Additional style options (color, width, etc.)

    Returns:
        pptx.shapes.Shape: The created connector shape
    """
    from pptx.enum.shapes import MSO_CONNECTOR

    start_x, start_y = start_pos
    end_x, end_y = end_pos

    # Create connector
    connector = slide.slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x),
        Inches(start_y),
        Inches(end_x),
        Inches(end_y),
    )

    # Line color
    if "color" in kwargs:
        connector.line.color.rgb = RGBColor.from_string(kwargs["color"])

    # Line width
    if "width" in kwargs:
        connector.line.width = Pt(kwargs["width"])

    # Arrow settings
    from pptx.enum.dml import MSO_LINE

    # Default to arrow at end
    connector.line.begin_style = MSO_LINE.NONE
    connector.line.end_style = MSO_LINE.ARROW

    # Set specific arrow styles if requested
    if "begin_arrow" in kwargs:
        if kwargs["begin_arrow"]:
            connector.line.begin_style = MSO_LINE.ARROW
        else:
            connector.line.begin_style = MSO_LINE.NONE

    if "end_arrow" in kwargs:
        if kwargs["end_arrow"]:
            connector.line.end_style = MSO_LINE.ARROW
        else:
            connector.line.end_style = MSO_LINE.NONE

    # Arrow size
    if "arrow_size" in kwargs:
        size = kwargs["arrow_size"]
        if size == "small":
            connector.line.begin_width = MSO_LINE.NARROW
            connector.line.end_width = MSO_LINE.NARROW
        elif size == "medium":
            connector.line.begin_width = MSO_LINE.MEDIUM
            connector.line.end_width = MSO_LINE.MEDIUM
        elif size == "large":
            connector.line.begin_width = MSO_LINE.WIDE
            connector.line.end_width = MSO_LINE.WIDE

    # Line dash style
    if "dash_style" in kwargs:
        dash_style = kwargs["dash_style"].lower()
        if dash_style == "solid":
            connector.line.dash_style = MSO_LINE.SOLID
        elif dash_style == "dash":
            connector.line.dash_style = MSO_LINE.DASH
        elif dash_style == "dot":
            connector.line.dash_style = MSO_LINE.ROUND_DOT
        elif dash_style == "dash_dot":
            connector.line.dash_style = MSO_LINE.DASH_DOT

    return connector


def add_image(slide, image_path, position, size=None, **kwargs):
    """
    Add an image to a slide.

    Args:
        slide (PPTXSlide): The slide object
        image_path (str): Path to the image file
        position (tuple): (x, y) position in inches
        size (tuple, optional): (width, height) size in inches. If None, use image's size.
        **kwargs: Additional options

    Returns:
        pptx.shapes.Picture: The created picture shape
    """
    x, y = position

    if size:
        width, height = size
        picture = slide.slide.shapes.add_picture(
            image_path, Inches(x), Inches(y), Inches(width), Inches(height)
        )
    else:
        picture = slide.slide.shapes.add_picture(image_path, Inches(x), Inches(y))

    # Apply any border if specified
    if "border_color" in kwargs:
        picture.line.color.rgb = RGBColor.from_string(kwargs["border_color"])

        # Line width
        if "border_width" in kwargs:
            picture.line.width = Pt(kwargs["border_width"])

    return picture
